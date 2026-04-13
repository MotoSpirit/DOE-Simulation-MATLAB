import sys
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

TEMPLATE = r'E:\MOTOSPIRIT\DOE\DOE_Final\PLANTILLA PROJECTES.pptx'
prs = Presentation(TEMPLATE)

print(f'Slides: {len(prs.slides)}')
print(f'Slide width:  {prs.slide_width.inches:.3f}"')
print(f'Slide height: {prs.slide_height.inches:.3f}"')
print()

for i, slide in enumerate(prs.slides):
    layout = slide.slide_layout
    print(f'--- Slide {i+1} (layout: "{layout.name}") ---')
    for shape in slide.shapes:
        t  = shape.top/914400   if shape.top   else 0
        l  = shape.left/914400  if shape.left  else 0
        w  = shape.width/914400 if shape.width  else 0
        h  = shape.height/914400 if shape.height else 0
        print(f'  [{shape.shape_type}] "{shape.name}"  pos=({l:.2f}",{t:.2f}")  size=({w:.2f}"x{h:.2f}")')
        if hasattr(shape, 'text') and shape.text.strip():
            print(f'     text: {repr(shape.text[:120])}')
        # Colors de farciment
        try:
            fill = shape.fill
            if fill.type == 1:  # SOLID
                rgb = fill.fore_color.rgb
                print(f'     fill: #{rgb}')
        except: pass
        if shape.shape_type == 13:
            print(f'     -> IMATGE/PICTURE')

print('\n--- Layouts disponibles ---')
for i, layout in enumerate(prs.slide_layouts):
    print(f'  [{i}] "{layout.name}"')

print('\n--- Slide Master colors principals ---')
master = prs.slide_masters[0]
for shape in master.shapes:
    try:
        fill = shape.fill
        if fill.type == 1:
            rgb = fill.fore_color.rgb
            t  = shape.top/914400   if shape.top   else 0
            l  = shape.left/914400  if shape.left  else 0
            w  = shape.width/914400 if shape.width  else 0
            h  = shape.height/914400 if shape.height else 0
            print(f'  "{shape.name}" fill=#{rgb}  pos=({l:.2f}",{t:.2f}")  size=({w:.2f}"x{h:.2f}")')
    except: pass
