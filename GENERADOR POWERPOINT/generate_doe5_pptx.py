"""
generate_doe5_pptx.py  - DOE 5 PowerPoint Generator
=====================================================
Genera un PowerPoint DOE 5 usant la PLANTILLA PROJECTES de Motospirit.
- Una sola slide per pestanya, amb totes les imatges individuals en graella
- Aspect-ratio preservat (no estira ni distorsiona)
- Les captures TAB no s'inclouen

Us:
    python generate_doe5_pptx.py <carpeta_png> [output.pptx]
"""

import sys, os, glob
from datetime import datetime

# Forcar UTF-8 a la consola de Windows
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
if hasattr(sys.stderr, 'reconfigure'):
    sys.stderr.reconfigure(encoding='utf-8', errors='replace')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from PIL import Image as PILImage
except ImportError as e:
    print(f"ERROR: {e}  ->  pip install python-pptx pillow")
    sys.exit(1)

# ── CONFIGURACIO ──────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE   = os.path.join(SCRIPT_DIR, "PLANTILLA PROJECTES.pptx")

SW_IN = 26.667   # slide width  (polzades)
SH_IN = 15.000   # slide height (polzades)

# Area de contingut (sota la capçalera de la plantilla)
CONTENT_LEFT   = 1.0
CONTENT_TOP    = 4.4
CONTENT_RIGHT  = SW_IN - 1.0
CONTENT_BOTTOM = SH_IN - 0.5
CONTENT_W      = CONTENT_RIGHT  - CONTENT_LEFT   # 24.667"
CONTENT_H      = CONTENT_BOTTOM - CONTENT_TOP    # 10.1"
GAP            = 0.15   # espai entre imatges (polzades)

TAB_ORDER = [
    ("01_BattPower",     "Battery Power"),
    ("02_EnginePower",   "Engine Power"),
    ("03_WheelPower",    "Wheel Power"),
    ("04_BatteryDeep",   "Battery Deep Analysis"),
    ("05_MotorAnalysis", "Motor Analysis"),
    ("06_ThermalMotor",  "Thermal Motor"),
    ("07_Derating",      "Derating Map"),
    ("08_VelocityProf",  "Velocity Profile (7 Laps)"),
    ("09_LapOverlay",    "Lap-to-Lap Overlay"),
    ("10_GearShifts",    "Gear Shifts"),
    ("11_TorqueLimits",  "Torque Limits"),
]

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x99, 0x99, 0x99)


# ── HELPERS ───────────────────────────────────────────────────────────────────

def I(x): return Inches(x)


def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name.lower() == name.lower():
            return layout
    return prs.slide_layouts[0]


def clear_slide_placeholders(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, italic=False,
                 color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return tb


def img_size_px(path):
    """Retorna (w_px, h_px) de la imatge."""
    try:
        with PILImage.open(path) as img:
            return img.size
    except Exception:
        return (1920, 1080)   # fallback


def fit_image_in_cell(img_w_px, img_h_px, cell_w_in, cell_h_in):
    """
    Calcula (off_x, off_y, draw_w, draw_h) en polzades per centrar la imatge
    dins la cel·la mantenint el ratio original (letterbox).
    """
    img_ratio  = img_w_px / img_h_px
    cell_ratio = cell_w_in / cell_h_in

    if img_ratio > cell_ratio:
        # mes ampla que la cel·la -> ajustar a l'amplada
        draw_w = cell_w_in
        draw_h = cell_w_in / img_ratio
    else:
        # mes alta que la cel·la -> ajustar a l'altura
        draw_h = cell_h_in
        draw_w = cell_h_in * img_ratio

    off_x = (cell_w_in - draw_w) / 2
    off_y = (cell_h_in - draw_h) / 2
    return off_x, off_y, draw_w, draw_h


def grid_layout(n):
    """
    Retorna (rows, cols) per a n imatges, prioritzant files menors.
    """
    if   n == 1: return (1, 1)
    elif n == 2: return (1, 2)
    elif n == 3: return (1, 3)
    elif n == 4: return (2, 2)
    elif n == 5: return (2, 3)   # una cel·la buida
    elif n == 6: return (2, 3)
    elif n == 7: return (2, 4)
    elif n == 8: return (2, 4)
    elif n == 9: return (3, 3)
    else:        return (3, 4)


def place_images_grid(slide, image_paths, area_left, area_top, area_w, area_h):
    """
    Col·loca una llista d'imatges en una graella dins l'area indicada,
    preservant el ratio de cada imatge.
    """
    n = len(image_paths)
    if n == 0:
        return

    rows, cols = grid_layout(n)

    cell_w = (area_w - GAP * (cols - 1)) / cols
    cell_h = (area_h - GAP * (rows - 1)) / rows

    for idx, path in enumerate(image_paths):
        row = idx // cols
        col = idx  % cols

        cell_left = area_left + col * (cell_w + GAP)
        cell_top  = area_top  + row * (cell_h + GAP)

        w_px, h_px = img_size_px(path)
        off_x, off_y, draw_w, draw_h = fit_image_in_cell(w_px, h_px, cell_w, cell_h)

        slide.shapes.add_picture(
            path,
            I(cell_left + off_x), I(cell_top + off_y),
            I(draw_w), I(draw_h)
        )


def find_images(folder, prefix):
    """Retorna (tab_png|None, [individual_pngs])."""
    tab_png  = glob.glob(os.path.join(folder, f"{prefix}_TAB.png"))
    ind_pngs = sorted(glob.glob(os.path.join(folder, f"{prefix}__*.png")))
    return (tab_png[0] if tab_png else None), ind_pngs


def read_params(folder):
    p = os.path.join(folder, "Parametres_Base_Simulacio.txt")
    if os.path.exists(p):
        with open(p, encoding='utf-8', errors='replace') as f:
            return f.read()
    return ""


# ── CAPÇALERA ESTANDARD ───────────────────────────────────────────────────────

def add_header(slide, title, subtitle=""):
    add_text_box(slide, "Dept. Electric  ·  DOE 5",
                 1.10, 0.84, 10.0, 0.55, font_size=16, color=WHITE)
    add_text_box(slide, datetime.now().strftime("%d.%m.%Y"),
                 22.57, 0.98, 3.0, 0.52, font_size=14,
                 color=GRAY, align=PP_ALIGN.RIGHT)
    add_text_box(slide, title.upper(),
                 1.34, 2.19, 20.0, 1.1, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     1.34, 3.45, 22.0, 0.7, font_size=19, color=GRAY)


# ── CREADORS DE SLIDE ─────────────────────────────────────────────────────────

def make_cover_slide(prs, folder_name, params_text):
    layout = get_layout(prs, "TITLE_AND_BODY")
    slide  = prs.slides.add_slide(layout)
    clear_slide_placeholders(slide)

    add_text_box(slide, "Departament Electric  ·  DOE 5",
                 1.10, 0.84, 12.0, 0.65, font_size=22, color=WHITE)
    add_text_box(slide, datetime.now().strftime("%d.%m.%Y"),
                 22.57, 1.24, 3.5, 0.52, font_size=18,
                 color=GRAY, align=PP_ALIGN.RIGHT)
    add_text_box(slide,
                 "GEARBOX vs DIRECT DRIVE\nPowertrain Performance Analysis  ·  DOE 5",
                 2.0, 5.5, 22.0, 3.5, font_size=52, bold=True, color=WHITE)
    add_text_box(slide, f"Simulacio: {folder_name}",
                 2.0, 9.2, 22.0, 0.8, font_size=22, color=GRAY)

    lines = [l for l in params_text.split('\n') if l.strip()][:18]
    add_text_box(slide, '\n'.join(lines),
                 2.0, 10.2, 22.5, 4.5, font_size=12, color=GRAY, italic=True)


def make_tab_slide(prs, tab_title, tab_png, ind_pngs):
    """
    Crea UNA sola slide per la pestanya:
    - Si hi ha individuals: graella amb tots els individuals (sense TAB)
    - Si nomes hi ha TAB (pestanya d'un sol axes): mostra el TAB
    """
    # Quines imatges mostrem?
    images_to_show = ind_pngs if ind_pngs else ([tab_png] if tab_png else [])
    if not images_to_show:
        return

    n = len(images_to_show)
    layout = get_layout(prs, "Light")
    slide  = prs.slides.add_slide(layout)
    clear_slide_placeholders(slide)

    subtitle = f"{n} grafic{'s' if n > 1 else ''}"
    add_header(slide, tab_title, subtitle)

    place_images_grid(
        slide, images_to_show,
        area_left = CONTENT_LEFT,
        area_top  = CONTENT_TOP,
        area_w    = CONTENT_W,
        area_h    = CONTENT_H
    )


def make_params_slide(prs, params_text):
    layout   = get_layout(prs, "Light")
    lines    = params_text.split('\n')
    chunk_sz = 38
    chunks   = [lines[i:i+chunk_sz] for i in range(0, len(lines), chunk_sz)]

    for ci, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(layout)
        clear_slide_placeholders(slide)
        suffix = f" ({ci+1}/{len(chunks)})" if len(chunks) > 1 else ""
        add_header(slide, f"Parametres Base de Simulacio{suffix}")
        add_text_box(slide, '\n'.join(chunk),
                     1.34, 4.5, 24.0, 10.0, font_size=11, color=GRAY)


def make_thankyou_slide(prs):
    layout = get_layout(prs, "Thank you")
    prs.slides.add_slide(layout)


# ── MAIN ──────────────────────────────────────────────────────────────────────

def generate(png_folder, output_path):
    if not os.path.isabs(png_folder):
        png_folder = os.path.join(SCRIPT_DIR, png_folder)

    if not os.path.isdir(png_folder):
        print(f"ERROR: No trobada la carpeta: {png_folder}")
        sys.exit(1)

    if not os.path.exists(TEMPLATE):
        print(f"ERROR: No trobada la plantilla: {TEMPLATE}")
        sys.exit(1)

    print(f"Plantilla:  {TEMPLATE}")
    print(f"Imatges de: {png_folder}")
    print(f"Sortida:    {output_path}")
    print()

    prs = Presentation(TEMPLATE)

    # Esborrar slides d'exemple de la plantilla
    xml_sldIdLst = prs.slides._sldIdLst
    for _ in range(len(prs.slides)):
        sldId = xml_sldIdLst[0]
        rId = sldId.get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        if rId:
            try: prs.part.drop_rel(rId)
            except: pass
        xml_sldIdLst.remove(sldId)

    params_text = read_params(png_folder)
    folder_name = os.path.basename(png_folder)

    # 1. Portada
    print("  -> Creant portada...")
    make_cover_slide(prs, folder_name, params_text)

    # 2. Una slide per pestanya
    for prefix, title in TAB_ORDER:
        tab_png, ind_pngs = find_images(png_folder, prefix)
        images = ind_pngs if ind_pngs else ([tab_png] if tab_png else [])
        if images:
            src = f"individuals: {len(ind_pngs)}" if ind_pngs else "TAB (1 axes)"
            print(f"  -> {title}  ({src})")
            make_tab_slide(prs, title, tab_png, ind_pngs)
        else:
            print(f"  -> {title}  (sense imatges, omes)")

    # 3. Parametres base
    if params_text:
        print("  -> Afegint parametres base...")
        make_params_slide(prs, params_text)

    # 4. Slide final
    make_thankyou_slide(prs)

    prs.save(output_path)
    print(f"\n[OK] PowerPoint generat: {output_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(0)

    png_folder = sys.argv[1]
    if len(sys.argv) >= 3:
        output_path = sys.argv[2]
    else:
        fn = os.path.basename(os.path.abspath(os.path.join(SCRIPT_DIR, png_folder)))
        output_path = os.path.join(SCRIPT_DIR, f"DOE5_{fn}.pptx")

    generate(png_folder, output_path)
